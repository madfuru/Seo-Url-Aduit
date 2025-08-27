import io, re, os, datetime as dt
import re as _re
import time
import json
from urllib.parse import urlparse
import logging
from concurrent.futures import ThreadPoolExecutor, as_completed, TimeoutError, wait
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from flask import Response
from werkzeug.exceptions import HTTPException
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx import Presentation
import os, requests
import threading

app = Flask(__name__)
# allow all origins (simplest); later you can restrict to your WP domain
CORS(app, resources={r"/*": {"origins": "*"}})

# ---------- helpers ----------
UA = {"User-Agent":"Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                   "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"}

def fetch(url, timeout=12):
    return requests.get(url, timeout=timeout, headers=UA, allow_redirects=True)

def safe_text(el):
    return (el.get_text(" ", strip=True) if el else "").strip()

def domain_from_url(u):
    return urlparse(u).netloc.lower()
TEMPLATE_PATH = os.getenv(
    "PPTX_TEMPLATE",
    os.path.join(os.path.dirname(__file__), "templates.pptx")  # Use the correct filename here
)

def _load_template():
    if not os.path.exists(TEMPLATE_PATH):
        raise FileNotFoundError(f"Template not found at {TEMPLATE_PATH}")
    try:
        return Presentation(TEMPLATE_PATH)
    except Exception as e:
        raise RuntimeError(f"Failed to load PPTX template: {e}")
# --- RDAP with quick universal endpoint first ---
def _rdap_events_to_dates(j):
    created = expires = None
    for e in (j.get("events") or []):
        act = (e.get("eventAction") or "").lower()
        if act in ("registration","creation"): created = e.get("eventDate")
        if act in ("expiration","expiry","expire"): expires = e.get("eventDate")
    fmt = lambda d: None if not d else d.split("T")[0]
    return fmt(created), fmt(expires)


def get_domain_info(u):
    dom = domain_from_url(u)

    # 1) Simple universal fallback first
    try:
        j = requests.get(f"https://rdap.org/domain/{dom}", timeout=8).json()
        created, expires = _rdap_events_to_dates(j)
        expiry_days = None
        if expires:
            d0 = dt.datetime.utcnow().date()
            d1 = dt.datetime.fromisoformat(expires.replace("Z","")).date()
            expiry_days = (d1 - d0).days
        registrar = None
        # try to pull a human registrar name if present
        for ent in j.get("entities", []) or []:
            if ent.get("roles") and "registrar" in [r.lower() for r in ent["roles"]]:
                va = (ent.get("vcardArray") or [None, []])[1]
                # find org/name fields in vcard
                for item in va:
                    if item and len(item) >= 4 and item[0] in ("fn","org"):
                        registrar = item[3]
                        break
                if registrar: break
        return {"created": created, "expiry": expires, "days_to_expire": expiry_days, "registrar": registrar}
    except Exception:
        pass

    # 2) IANA bootstrap fallback
    try:
        tld = dom.split(".")[-1]
        boot = requests.get("https://data.iana.org/rdap/dns.json", timeout=8).json()
        servers = [s for s in boot.get("services", []) if any(tld == x for x in s[0])]
        if servers:
            base = servers[0][1][0].rstrip("/")
            j = requests.get(f"{base}/domain/{dom}", timeout=8).json()
            created, expires = _rdap_events_to_dates(j)
            expiry_days = None
            if expires:
                d0 = dt.datetime.utcnow().date()
                d1 = dt.datetime.fromisoformat(expires.replace("Z","")).date()
                expiry_days = (d1 - d0).days
            registrar = None
            for ent in j.get("entities", []) or []:
                if ent.get("roles") and "registrar" in [r.lower() for r in ent["roles"]]:
                    va = (ent.get("vcardArray") or [None, []])[1]
                    for item in va:
                        if item and len(item) >= 4 and item[0] in ("fn","org"):
                            registrar = item[3]
                            break
                    if registrar: break
            return {"created": created, "expiry": expires, "days_to_expire": expiry_days, "registrar": registrar}
    except Exception:
        pass

    return {"created": None, "expiry": None, "days_to_expire": None, "registrar": None}

def classify_issue(raw):
    msg = raw.get("message","").lower()
    sev = (raw.get("severity") or "").lower()

    # Default
    name, itype, prio = raw.get("message",""), "Issue", "Medium"

    # H1 issues
    if "missing h1" in msg:
        name, itype, prio = "H1: Missing", "Issue", "Medium"
    elif "multiple h1" in msg:
        name, itype, prio = "H1: Multiple", "Warning", "Medium"

    # Images
    elif "without alt" in msg:
        name, itype, prio = "Images: Missing ALT", "Issue", "Medium"
    elif "large image" in msg:
        name, itype, prio = "Images: Over 100 KB", "Opportunity", "Medium"

    # Canonical
    elif "missing canonical" in msg:
        name, itype, prio = "Canonicals: Missing", "Issue", "High"
    elif "canonical url not returning" in msg or "not reachable" in msg:
        name, itype, prio = "Canonicals: Canonicalised", "Warning", "High"

    # Titles / Descriptions
    elif "missing <title>" in msg:
        name, itype, prio = "Page Titles: Missing", "Issue", "High"
    elif "title length > 60" in msg:
        name, itype, prio = "Page Titles: Over 60 Characters", "Opportunity", "Medium"
    elif "meta description > 160" in msg:
        name, itype, prio = "Meta Descriptions: Over 160 Characters", "Opportunity", "Medium"
    elif "missing meta description" in msg:
        name, itype, prio = "Meta Descriptions: Missing", "Warning", "High"

    # Robots/noindex
    elif "noindex" in msg:
        name, itype, prio = "Response Codes: Noindex Tag", "Issue", "High"

    # HTTP errors
    elif "http status" in msg:
        name, itype, prio = "Response Codes: Error", "Issue", "High"

    return {"issue_name": name, "issue_type": itype, "priority": prio}



# ---------- PSI ----------
PSI_SEM = threading.Semaphore(2)  # limit concurrent /psi calls per worker
PSI_CACHE = {}                    # (url, strategy) -> (timestamp, json)
PSI_TTL = 60*60*24                # 24h
PSI_CACHE_MAX = 50  

# at top of file (once)
from math import ceil
def _psi_call(url, strategy, key=None, timeout=60, tries=2):
    # ensure scheme
    if not url.lower().startswith(("http://", "https://")):
        url = "https://" + url

    now = time.time()
    ck = (url, strategy)
    cached = PSI_CACHE.get(ck)
    if cached and now - cached[0] < PSI_TTL:
        return cached[1]

    base = "https://www.googleapis.com/pagespeedonline/v5/runPagespeed"

    def do():
        params = {"url": url, "strategy": strategy}
        if key:
            params["key"] = key
        return requests.get(base, params=params, timeout=timeout)

    last_err = None
    for attempt in range(1, tries+1):
        try:
            r = do()
            # Try to parse JSON always
            try:
                data = r.json()
            except Exception:
                data = {"error": {"message": f"Non-JSON from PSI (HTTP {r.status_code})"}}

            # If PSI returned error, include status/code details for visibility
            if "error" in data:
                err = data["error"]
                # Example: {"code":400,"status":"INVALID_ARGUMENT","message":"..."}
                msg = err.get("message","PSI error")
                status = err.get("status")
                code = err.get("code")
                data["error"]["message"] = f"{msg} (status={status}, code={code})"

            # cache (with eviction)
            if len(PSI_CACHE) >= PSI_CACHE_MAX:
                oldest = min(PSI_CACHE, key=lambda k: PSI_CACHE[k][0])
                PSI_CACHE.pop(oldest, None)
            PSI_CACHE[ck] = (time.time(), data)
            return data

        except (requests.Timeout, requests.ConnectionError) as e:
            last_err = f"{type(e).__name__}: {e}"
            time.sleep(0.7 * (2 ** (attempt-1)))  # small backoff
        except Exception as e:
            last_err = f"{type(e).__name__}: {e}"
            break

    return {"error": {"message": f"PSI request failed after {tries} tries: {last_err}"}}


def _psi_parse(j):
    if not j or "error" in j or not j.get("lighthouseResult"):
        # return the error string so UI can show it
        return None, (j.get("error", {}) or {}).get("message", "No PSI data")

    lhr = j["lighthouseResult"]
    audits = lhr.get("audits", {})
    cats = lhr.get("categories", {})
    perf = cats.get("performance", {}).get("score")

    def val(k):
        a = audits.get(k, {}) or {}
        return a.get("numericValue", a.get("displayValue"))

    return {
        "performance_score": round((perf or 0) * 100) if perf is not None else None,
        "fcp": val("first-contentful-paint"),
        "lcp": val("largest-contentful-paint"),
        "cls": (audits.get("cumulative-layout-shift", {}) or {}).get("numericValue"),
        "tbt": (audits.get("total-blocking-time", {}) or {}).get("numericValue"),
        "inp": (audits.get("experimental-interaction-to-next-paint", {}) or {}).get("numericValue"),
    }, None




def get_pagespeed_both(url):
    key = os.getenv("PSI_KEY")  # must be ENV VAR NAME
    # Desktop first
    d_raw = _psi_call(url, "desktop", key=key, timeout=60, tries=2)
    desktop, d_err = _psi_parse(d_raw)
    # Mobile next (shorter on purpose)
    m_raw = _psi_call(url, "mobile",  key=key, timeout=60, tries=1)
    mobile,  m_err = _psi_parse(m_raw)
    return {
        "desktop": desktop, "desktop_error": d_err,
        "mobile":  mobile,  "mobile_error":  m_err,
        "using_key": bool(key)
    }



# ---------- fast link checker ----------
def _fast_status(u):
    try:
        r = requests.head(u, timeout=3, allow_redirects=True, headers=UA)
        code = r.status_code
        if code in (403, 405):
            r = requests.get(u, timeout=4, allow_redirects=True, stream=False, headers=UA)
            code = r.status_code
        return {"href": u, "status": code}
    except Exception:
        return {"href": u, "status": 0}
def classify_issue(raw):
    """Map raw analyzer messages to: Issue Name / Issue Type / Priority."""
    msg = (raw.get("message") or "").lower()

    # defaults
    name, itype, prio = raw.get("message",""), "Issue", "Medium"

    # H1
    if "missing h1" in msg:
        return {"issue_name": "H1: Missing", "issue_type": "Issue", "priority": "Medium"}
    if "multiple h1" in msg:
        return {"issue_name": "H1: Multiple", "issue_type": "Warning", "priority": "Medium"}

    # Titles
    if "missing <title>" in msg:
        return {"issue_name": "Page Titles: Missing", "issue_type": "Issue", "priority": "High"}
    if "title length > 60" in msg or "over 60 characters" in msg:
        return {"issue_name": "Page Titles: Over 60 Characters", "issue_type": "Opportunity", "priority": "Medium"}
    if "below 30" in msg:
        return {"issue_name": "Page Titles: Below 30 Characters", "issue_type": "Opportunity", "priority": "Medium"}

    # Meta descriptions
    if "missing meta description" in msg:
        return {"issue_name": "Meta Descriptions: Missing", "issue_type": "Warning", "priority": "High"}
    if "meta description > 160" in msg or "over 160" in msg:
        return {"issue_name": "Meta Descriptions: Over 160 Characters", "issue_type": "Opportunity", "priority": "Medium"}

    # Canonical
    if "missing canonical" in msg:
        return {"issue_name": "Canonicals: Missing", "issue_type": "Issue", "priority": "High"}
    if "canonical url not returning" in msg or "canonical url not reachable" in msg:
        return {"issue_name": "Canonicals: Canonicalised", "issue_type": "Warning", "priority": "High"}
    if "non-indexable canonical" in msg:
        return {"issue_name": "Canonicals: Non-Indexable Canonical", "issue_type": "Issue", "priority": "High"}

    # Robots / noindex
    if "noindex" in msg:
        return {"issue_name": "Indexability: noindex Present", "issue_type": "Issue", "priority": "High"}

    # Images
    if "image(s) without alt" in msg:
        return {"issue_name": "Images: Missing ALT", "issue_type": "Issue", "priority": "Medium"}
    if "large image" in msg or " >150kb" in msg or "over 100 kb" in msg:
        return {"issue_name": "Images: Over 100 KB", "issue_type": "Opportunity", "priority": "Medium"}

    # HTTP / response
    if "http status" in msg:
        # you already add exact status text; treat 4xx/5xx as High
        return {"issue_name": "Response Codes: Error", "issue_type": "Issue", "priority": "High"}

    # Canonical / robots fallback
    if "robots.txt" in msg:
        return {"issue_name": "Response Codes: Blocked by robots.txt", "issue_type": "Warning", "priority": "High"}

    # Fallback
    return {"issue_name": raw.get("message","").strip() or "Unknown", "issue_type": "Issue", "priority": "Medium"}


# ---------- analyzer ----------
def analyze_html(url):
    issues, links = [], []
    try:
        res = fetch(url, timeout=8)
        if res.status_code != 200:
            issues.append({"severity":"error","message":f"HTTP status {res.status_code} (not 200)."})
        if not url.startswith("https://"):
            issues.append({"severity":"warn","message":"URL is not HTTPS."})

        soup = BeautifulSoup(res.text, "html.parser")

        title = (soup.find("title").get_text(" ", strip=True) if soup.find("title") else "").strip()
        if not title: issues.append({"severity":"error","message":"Missing <title>."})
        elif len(title) > 60: issues.append({"severity":"warn","message":"Title length > 60 characters."})

        md = soup.find("meta", attrs={"name":"description"})
        desc = (md.get("content","").strip() if md else "")
        if not desc: issues.append({"severity":"warn","message":"Missing meta description."})
        elif len(desc) > 160: issues.append({"severity":"warn","message":"Meta description > 160 characters."})

        h1s = [h.get_text(" ", strip=True) for h in soup.find_all(_re.compile("^h1$"))]
        if len(h1s)==0: issues.append({"severity":"warn","message":"Missing H1."})
        elif len(h1s)>1: issues.append({"severity":"warn","message":f"Multiple H1s ({len(h1s)})."})

        can = soup.find("link", rel="canonical")
        if not can or not can.get("href"):
            issues.append({"severity":"warn","message":"Missing canonical link."})
        else:
            try:
                cr = fetch(can.get("href"), timeout=4)
                if cr.status_code != 200:
                    issues.append({"severity":"warn","message":"Canonical URL not returning 200."})
            except Exception:
                issues.append({"severity":"warn","message":"Canonical URL not reachable."})

        robots_meta = soup.find("meta", attrs={"name":"robots"})
        if robots_meta and "noindex" in robots_meta.get("content","").lower():
            issues.append({"severity":"error","message":"Page has noindex meta."})

        # super-fast image checks (no big HEAD calls)
        noalt = sum(1 for img in soup.find_all("img") if not (img.get("alt") or "").strip())
        if noalt > 0:
            issues.append({"severity":"warn","message":f"{noalt} image(s) without alt."})

        # tiny same-domain link sample (<=6) with short timeouts
        page_host = urlparse(url).netloc.lower()
        cands = []
        for a in soup.find_all("a", limit=120):
            href = a.get("href"); 
            if not href or href.startswith("#"): continue
            full = href if href.startswith("http") else requests.compat.urljoin(url, href)
            if urlparse(full).netloc.lower() == page_host:
                cands.append(full)
            if len(cands) >= 6: break

        if cands:
            with ThreadPoolExecutor(max_workers=6) as ex:
                futs = [ex.submit(lambda u: requests.head(u, timeout=2, allow_redirects=True, headers=UA), c) for c in cands]
                for i, f in enumerate(futs):
                    try:
                        r = f.result(timeout=3)
                        links.append({"href": cands[i], "status": r.status_code})
                    except Exception:
                        links.append({"href": cands[i], "status": 0})

        return {"issues": issues, "links": links}
    except Exception:
        return {"issues":[{"severity":"error","message":"Fetch failed or invalid HTML."}], "links":[]}
# ---------- routes ----------
@app.route("/")
def health():
    return jsonify({"ok": True, "service": "seo-url-audit"})

@app.route("/psi")
def psi():
    url = request.args.get("url","").strip()
    if not url: return jsonify({"error":"missing url"}), 400
    if not os.getenv("PSI_KEY"): return jsonify({"error":"psikey_missing"}), 200
    if not PSI_SEM.acquire(blocking=False):
        return jsonify({"error":"busy", "message":"Please retry in a moment"}), 429
    try:
        return jsonify(get_pagespeed_both(url)), 200
    finally:
        PSI_SEM.release()

@app.errorhandler(Exception)
def _any_error(e):
    # Log the stack trace in Render logs, but don't 500 the browser.
    app.logger.exception(f"Unhandled error: {e}")
    return jsonify({"error": "server_error", "detail": str(e)[:120]}), 200

@app.route("/debug/psikey")
def debug_psikey():
    import os
    return jsonify({"using_key": bool(os.getenv("PSI_KEY"))})

@app.route("/analyze")
def analyze():
    url = request.args.get("url","").strip()
    if not url:
        return jsonify({"error":"missing url"}), 400

    skip_psi = request.args.get("skip_psi", "0") == "1"

    results = {"speed": {}, "domain": {}, "seo": {"issues": [], "classified": []}, "links": []}

    def _speed():  return ("speed", get_pagespeed_both(url))
    def _domain(): return ("domain", get_domain_info(url))
    def _seo():
        x = analyze_html(url)
        return ("seo_links", x)

    # run in parallel, but optionally skip PSI for a fast response
    with ThreadPoolExecutor(max_workers=3) as ex:
        futs = []
        if not skip_psi:
            futs.append(ex.submit(_speed))
        futs += [ex.submit(_domain), ex.submit(_seo)]

        done, not_done = wait(futs, timeout=14)  # keep fast for UI
        for f in not_done:
            f.cancel()

        for f in done:
            try:
                k, v = f.result()
                if k == "seo_links":
                    issues = v.get("issues", [])
                    results["seo"] = {
                        "issues": issues,
                        "classified": [classify_issue(i) for i in issues]
                    }
                    results["links"] = v.get("links", [])
                else:
                    results[k] = v
            except Exception as e:
                app.logger.warning(f"worker {k if 'k' in locals() else '?'} failed: {e}")

    return jsonify(results), 200


@app.route("/favicon.ico")
def favicon():
    return Response(status=204)

@app.route("/debug/template")
def debug_template():
    import os
    exists = os.path.exists(TEMPLATE_PATH)
    return jsonify({"template_path": TEMPLATE_PATH, "exists": exists})

@app.route("/report")
def report():
    url = request.args.get("url","").strip()
    prs = Presentation("/path/to/templates.pptx")  # local path
    if not url:
        return "missing url", 400

    # Gather data
    speed_both = get_pagespeed_both(url)  # {desktop, desktop_error, mobile, mobile_error}
    domain = get_domain_info(url)
    seo = analyze_html(url)
    issues = seo.get("issues", [])
    classified = [classify_issue(i) for i in issues]

    # Flatten speed for placeholders
    dsk, mob = speed_both.get("desktop") or {}, speed_both.get("mobile") or {}

    replacements = {
        "{{URL}}": url,
        "{{Created}}": domain.get("created"),
        "{{Expiry}}": domain.get("expiry"),
        "{{DaysToExpire}}": domain.get("days_to_expire"),
        "{{Registrar}}": domain.get("registrar"),

        "{{PerfDesktop}}": dsk.get("performance_score"),
        "{{FCPDesktop}}": dsk.get("fcp"),
        "{{LCPDesktop}}": dsk.get("lcp"),
        "{{CLSDesktop}}": dsk.get("cls"),
        "{{INPDesktop}}": dsk.get("inp") or dsk.get("tbt"),

        "{{PerfMobile}}": mob.get("performance_score"),
        "{{FCPMobile}}": mob.get("fcp"),
        "{{LCPMobile}}": mob.get("lcp"),
        "{{CLSMobile}}": mob.get("cls"),
        "{{INPMobile}}": mob.get("inp") or mob.get("tbt"),
    }

    # Replace placeholders in all slides
    replace_text_everywhere(prs, replacements)

    # Insert the On Page Analysis Overview table (find slide by title text)
    slide = find_slide_by_title(prs, "On Page Analysis Overview")
    if slide:
        add_overview_table(slide, classified[:25])  # cap to keep layout tidy

    # Stream back
    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    fname = f"SEO_Audit_{domain_from_url(url)}.pptx"
    return send_file(
        buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=fname
    )

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000)

